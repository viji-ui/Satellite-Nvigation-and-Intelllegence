import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # -------------------------------------------------------------
    # Helper constraints
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]

    # Set background color or theme preferences (Optional, sticking to clean defaults)
    
    # -------------------------------------------------------------
    # Slide 1: Title
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Satellite Navigation Overview"
    subtitle.text = "NLP Pipeline Architecture\nTransforming raw satellite data into actionable insights"

    # -------------------------------------------------------------
    # Slide 2: Pipeline Overview
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Overview of the NLP Pipeline"
    tf = body.text_frame
    tf.text = "The pipeline processes 2,000+ raw satellite records (UCS Database)."
    
    p = tf.add_paragraph()
    p.text = "Extracts rich insights through 5 sequential steps."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Core Goal: Convert raw statistics into human-readable text, analyze sentiments, evaluate capabilities, and find natural satellite clusters."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "No heavy LLMs required — uses highly optimized, rule-based and classical ML techniques (VADER, TF-IDF, K-Means)."
    p.level = 0

    # -------------------------------------------------------------
    # Slide 3: Step 1 & 2 - Data Prep & Text Generation
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Steps 1 & 2: Data Cleaning & Text Generation"
    tf = body.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Data Cleaning (data_cleaner.py):"
    
    p = tf.add_paragraph()
    p.text = "Loads raw CSV, filters necessary columns, handles missing values, and prepares orbital/launch data."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Text Generation Engine (text_generator.py):"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "A deterministic template system that converts raw tabular data (e.g., orbit class, mass, purpose) into a human-readable mission profile."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Outcome: Rich, standardized text descriptions generated for every single satellite."
    p.level = 1

    # -------------------------------------------------------------
    # Slide 4: Step 3 - Sentiment Analysis
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Step 3: Sentiment Analysis"
    tf = body.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Method: VADER-style rule-based sentiment scoring."
    
    p = tf.add_paragraph()
    p.text = "Purpose: Analyzes the newly generated mission descriptions to detect the operational \"tone\" (Positive, Neutral, Negative)."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Key Insight:"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Civilian / LEO Earth observation missions generally show positive or neutral tones."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Military-focused missions lean heavily towards a cautious, strict, or neutral operational tone."
    p.level = 1

    # -------------------------------------------------------------
    # Slide 5: Step 4 - Capability Scoring
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Step 4: Strategic Capability Scoring"
    tf = body.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Method: Rule-based keyword extraction and mission assessment."
    
    p = tf.add_paragraph()
    p.text = "Purpose: Evaluates each satellite's strategic versus civilian capability."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Process:"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Examines mission intent, operational users (e.g., Military vs. Civil), and orbit logistics."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Calculates a 'Capability Score' useful for identifying critical defense or communications infrastructure."
    p.level = 1

    # -------------------------------------------------------------
    # Slide 6: Step 5 - Machine Learning Clustering
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Step 5: NLP Clustering & Archetypes"
    tf = body.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Method: TF-IDF (Term Frequency-Inverse Document Frequency) paired with K-Means Clustering."
    
    p = tf.add_paragraph()
    p.text = "Purpose: Automatically groups the satellite fleet into 6 distinct behavioral archetypes based on their mission text."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Results: Segments the thousands of individual satellites into natural groups:"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "E.g., \"Earth Observation LEO\", \"Strategic Military Comms\", or \"Deep Space/Science\"."
    p.level = 1

    # -------------------------------------------------------------
    # Slide 7: Integrated Outputs
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Pipeline Outputs & Application Integration"
    tf = body.text_frame
    
    p = tf.paragraphs[0]
    p.text = "The NLP pipeline directly feeds the interactive web dashboard."
    
    p = tf.add_paragraph()
    p.text = "satellite_profiles.csv: The full, enriched tabular dataset with all newly generated NLP columns."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "satellite_data.json: Lightweight JSON format consumed by the frontend."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Visualized In Dashboard:"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Dynamic Word Clouds based on Mission Texts."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Sentiment by Orbit distributions."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Strategic Capability visualizers."
    p.level = 1

    # Save presentation
    prs.save("NLP_Pipeline_Presentation.pptx")
    print("Successfully created NLP_Pipeline_Presentation.pptx!")

if __name__ == "__main__":
    create_presentation()
